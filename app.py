import os
import json
from datetime import datetime

from flask import Flask, render_template, request, send_file
from werkzeug.utils import secure_filename
from pypdf import PdfReader
from dotenv import load_dotenv
from google import genai
from pptx import Presentation

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet


load_dotenv()

app = Flask(__name__, template_folder="templates")

app.config["MAX_CONTENT_LENGTH"] = 10 * 1024 * 1024

UPLOAD_FOLDER = os.path.join(os.getcwd(), "upload")
GENERATED_FOLDER = os.path.join(os.getcwd(), "generated")
DADOS_FOLDER = os.path.join(os.getcwd(), "dados")

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GENERATED_FOLDER, exist_ok=True)
os.makedirs(DADOS_FOLDER, exist_ok=True)

client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))


@app.route("/")
def index():
    return render_template("index.html")


def extrair_texto_pdf(caminho_pdf):
    reader = PdfReader(caminho_pdf)

    texto = ""

    for pagina in reader.pages:
        texto += pagina.extract_text() or ""

    return texto


def gerar_perguntas_pdf(texto, quantidade, dificuldade, alternativas):
    prompt = f"""
Você é um professor especialista em elaboração de questões educacionais.

Com base no conteúdo abaixo, crie {quantidade} questões de múltipla escolha de nível {dificuldade}.

Regras obrigatórias:
Não use caracteres especiais como *, _, #, > ou formatação Markdown.
As questões devem ser bem contextualizadas.
Evite perguntas óbvias ou que apenas copiem trechos do texto.
Cada questão deve ter uma pequena contextualização antes da pergunta.
As alternativas devem ser curtas e objetivas.
Crie {alternativas} alternativas por questão.
Apenas uma alternativa deve estar correta.
Não coloque a resposta correta logo após cada questão.
Crie o gabarito somente no final.
No gabarito, informe apenas o número da questão e a letra correta.
Use linguagem adequada para estudantes do ensino médio.
Não faça menção ao texto base nas questões.

Formato esperado:

QUESTÃO 1
Contextualização da questão.

Pergunta?

A) Alternativa
B) Alternativa
C) Alternativa
D) Alternativa

GABARITO

1 - B
2 - D

Conteúdo para análise:
{texto}
"""

    resposta = client.models.generate_content(
        model="gemini-3.5-flash",
        contents=prompt
    )

    return resposta.text


def gerar_perguntas_quiz(texto, quantidade, dificuldade, alternativas):
    prompt = f"""
Gere apenas um array JSON válido.
A resposta deve começar com [ e terminar com ].
Não escreva nenhuma frase antes ou depois.
Não use markdown.
Não faça qualquer menção ao texto base.

Crie {quantidade} perguntas de múltipla escolha de nível {dificuldade}.
Cada pergunta deve ter exatamente {alternativas} alternativas.

Formato:
[
  {{
    "pergunta": "Texto da pergunta",
    "alternativas": ["Alternativa 1", "Alternativa 2", "Alternativa 3", "Alternativa 4"],
    "resposta": "Alternativa correta"
  }}
]

Conteúdo:
{texto}
"""

    resposta = client.models.generate_content(
        model="gemini-3.5-flash",
        contents=prompt
    )

    texto_resposta = resposta.text.strip()
    texto_resposta = texto_resposta.replace("```json", "")
    texto_resposta = texto_resposta.replace("```", "")
    texto_resposta = texto_resposta.strip()

    inicio = texto_resposta.find("[")
    fim = texto_resposta.rfind("]") + 1

    if inicio == -1 or fim == 0:
        print("Resposta da IA:", texto_resposta)
        raise ValueError("A IA não retornou uma lista JSON válida.")

    json_limpo = texto_resposta[inicio:fim]

    print("JSON LIMPO:", json_limpo)

    return json.loads(json_limpo)


def salvar_resultado_pdf(perguntas):
    caminho = os.path.join(DADOS_FOLDER, "perguntas_pdf.json")

    registro = {
        "data": datetime.now().strftime("%d/%m/%Y %H:%M:%S"),
        "perguntas": perguntas
    }

    if os.path.exists(caminho):
        with open(caminho, "r", encoding="utf-8") as arquivo:
            dados = json.load(arquivo)
    else:
        dados = []

    dados.append(registro)

    with open(caminho, "w", encoding="utf-8") as arquivo:
        json.dump(dados, arquivo, ensure_ascii=False, indent=4)

def extrair_texto_pptx(caminho_pptx):
    apresentacao = Presentation(caminho_pptx)

    texto = ""

    for slide in apresentacao.slides:
        for elemento in slide.shapes:
            if hasattr(elemento, "text"):
                texto += elemento.text + "\n"

    return texto


def salvar_quiz(perguntas):
    caminho = os.path.join(DADOS_FOLDER, "perguntas_quiz.json")

    with open(caminho, "w", encoding="utf-8") as arquivo:
        json.dump(perguntas, arquivo, ensure_ascii=False, indent=4)


def gerar_pdf(perguntas):
    caminho_pdf = os.path.join(GENERATED_FOLDER, "questoes.pdf")

    doc = SimpleDocTemplate(caminho_pdf)
    estilos = getSampleStyleSheet()

    conteudo = []

    conteudo.append(
        Paragraph(
            "Questões Geradas pelo LearnGen.IA",
            estilos["Title"]
        )
    )

    conteudo.append(Spacer(1, 20))

    texto_formatado = perguntas.replace("\n", "<br/>")

    conteudo.append(
        Paragraph(
            texto_formatado,
            estilos["BodyText"]
        )
    )

    doc.build(conteudo)

    return caminho_pdf


@app.route("/upload", methods=["POST"])
def upload():
    arquivo = request.files["pdf"]

    if arquivo.filename == "":
        return render_template(
            "index.html",
            mensagem="Selecione um PDF antes de enviar."
        )

    quantidade = request.form["quantidade"]
    dificuldade = request.form["dificuldade"]
    alternativas = request.form["alternativas"]
    tipo_saida = request.form["tipo_saida"]

    caminho = os.path.join(
        UPLOAD_FOLDER,
        secure_filename(arquivo.filename)
    )

    arquivo.save(caminho)

    nome_arquivo = arquivo.filename.lower()

    if nome_arquivo.endswith(".pdf"):
        texto = extrair_texto_pdf(caminho)

    elif nome_arquivo.endswith(".pptx"):
        texto = extrair_texto_pptx(caminho)

    else:
        return render_template(
            "index.html",
            mensagem="Formato inválido. Envie apenas PDF ou PPTX."
        )

    if tipo_saida == "pdf":
        perguntas = gerar_perguntas_pdf(
            texto,
            quantidade,
            dificuldade,
            alternativas
        )

        salvar_resultado_pdf(perguntas)
        gerar_pdf(perguntas)

        return render_template(
            "resultado.html",
            perguntas=perguntas
        )

    if tipo_saida == "quiz":
        perguntas = gerar_perguntas_quiz(
            texto,
            quantidade,
            dificuldade,
            alternativas
        )

        salvar_quiz(perguntas)

        return render_template("quiz.html")

    return "Tipo de saída inválido."


@app.route("/download")
def download():
    caminho_pdf = os.path.join(GENERATED_FOLDER, "questoes.pdf")

    return send_file(
        caminho_pdf,
        as_attachment=True
    )


@app.route("/api/perguntas")
def api_perguntas():

    with open(
        "dados/perguntas_quiz.json",
        "r",
        encoding="utf-8"
    ) as arquivo:

        perguntas = json.load(arquivo)

    return perguntas


@app.route("/novo_quiz")
def novo_quiz():

    with open(
        "dados/perguntas_quiz.json",
        "r",
        encoding="utf-8"
    ) as arquivo:

        perguntas = json.load(arquivo)

    return render_template("quiz.html")

@app.route("/sobre")
def sobre():
    return render_template("sobre.html")


if __name__ == "__main__":
    app.run(debug=True)